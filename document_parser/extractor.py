import io
import pypdf
import docx
import pptx
import openpyxl
import pandas as pd
from PIL import Image
import os

def extract_text_from_file(file_input):
    """
    Unified extractor for various file types.
    params: file_input (Streamlit UploadedFile OR str/Path for local file)
    returns: str (Extracted Text)
    """
    # Determine filename and file object
    if isinstance(file_input, str):
        filename = os.path.basename(file_input)
        file_obj = open(file_input, 'rb')
        should_close = True
    else:
        # Streamlit UploadedFile
        filename = file_input.name
        file_obj = file_input
        should_close = False
        
    file_type = filename.split('.')[-1].lower()
    
    try:
        if file_type == 'pdf':
            text = _extract_pdf(file_obj)
        elif file_type in ['docx', 'doc']:
            text = _extract_docx(file_obj)
        elif file_type in ['pptx', 'ppt']:
            text = _extract_pptx(file_obj)
        elif file_type in ['xlsx', 'xls']:
            text = _extract_excel(file_obj)
        elif file_type == 'csv':
             df = pd.read_csv(file_obj)
             text = df.to_string()
        elif file_type in ['txt', 'md', 'py', 'json']:
            # Reset pointer for text reading if needed, though usually at 0
            if hasattr(file_obj, 'seek'): file_obj.seek(0)
            text = str(file_obj.read().decode("utf-8"))
        elif file_type in ['png', 'jpg', 'jpeg']:
            text = "Valid Image File" 
        else:
            text = f"Unsupported file type: {file_type}"
            
        return text
        
    except Exception as e:
        return f"Error extracting file {filename}: {str(e)}"
    finally:
        if should_close:
            file_obj.close()

def _extract_pdf(file):
    pdf = pypdf.PdfReader(file)
    text = ""
    for page in pdf.pages:
        text += page.extract_text() + "\n"
    return text

def _extract_docx(file):
    doc = docx.Document(file)
    return "\n".join([p.text for p in doc.paragraphs])

def _extract_pptx(file):
    prs = pptx.Presentation(file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def _extract_excel(file):
    # Convert to string representation
    df = pd.read_excel(file)
    return df.to_string()
